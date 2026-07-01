"""Generate a PowerPoint deck from a structured spec.

Usage:
    python generate.py <spec.json|spec.yaml> [--output out.pptx] [--template tmpl.pptx]

Spec schema (JSON or YAML):
    title: str                      # deck title (used if no title slide given)
    subtitle: str                   # optional
    author: str                     # optional; written to notes
    slides: list[dict]              # each dict has a "type" key

Slide types:
    title       -> {title, subtitle}
    section     -> {title}
    bullets     -> {title, bullets: list[str]}
    two_column  -> {title, left_title, left_bullets, right_title, right_bullets}
    table       -> {title, headers: list[str], rows: list[list[str]]}
    image       -> {title, image_path: str}

Each slide may also carry `notes: str` which becomes speaker notes.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Inches, Pt


SLIDE_LAYOUTS = {
    # Default template layout indices (python-pptx default master).
    "title": 0,
    "section": 2,
    "bullets": 1,
    "two_column": 3,
    "table": 5,
    "image": 5,
}


def load_spec(path: Path) -> dict[str, Any]:
    text = path.read_text()
    if path.suffix.lower() in {".yaml", ".yml"}:
        import yaml
        return yaml.safe_load(text)
    return json.loads(text)


def build_deck(spec: dict[str, Any], template: Path | None = None) -> Presentation:
    prs = Presentation(str(template)) if template else Presentation()

    slides = spec.get("slides") or []
    if not slides:
        # No slides given — fall back to a single title slide from top-level fields.
        slides = [{
            "type": "title",
            "title": spec.get("title", "Untitled"),
            "subtitle": spec.get("subtitle", ""),
        }]

    for s in slides:
        kind = s.get("type", "bullets")
        builder = _builders.get(kind)
        if builder is None:
            raise ValueError(f"Unknown slide type: {kind!r}")
        builder(prs, s)

    return prs


def _title_slide(prs: Presentation, s: dict[str, Any]) -> None:
    layout = prs.slide_layouts[SLIDE_LAYOUTS["title"]]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s.get("title", "")
    subtitle = _find_placeholder(slide, 1)
    if subtitle is not None:
        subtitle.text = s.get("subtitle", "")
    _apply_notes(slide, s.get("notes"))


def _section_slide(prs: Presentation, s: dict[str, Any]) -> None:
    layout = prs.slide_layouts[SLIDE_LAYOUTS["section"]]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s.get("title", "")
    _apply_notes(slide, s.get("notes"))


def _bullets_slide(prs: Presentation, s: dict[str, Any]) -> None:
    layout = prs.slide_layouts[SLIDE_LAYOUTS["bullets"]]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s.get("title", "")
    body = _find_placeholder(slide, 1)
    if body is not None:
        tf = body.text_frame
        tf.clear()
        bullets = s.get("bullets", []) or [""]
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(line)
            para.level = 0
    _apply_notes(slide, s.get("notes"))


def _two_column_slide(prs: Presentation, s: dict[str, Any]) -> None:
    # python-pptx default master doesn't always have a two-content layout at
    # a consistent index; fall back to bullets + manually add a second box.
    try:
        layout = prs.slide_layouts[SLIDE_LAYOUTS["two_column"]]
    except IndexError:
        layout = prs.slide_layouts[SLIDE_LAYOUTS["bullets"]]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s.get("title", "")

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    top = Inches(1.75)
    height = slide_h - top - Inches(0.5)
    col_w = (slide_w - Inches(1.0)) // 2

    left = slide.shapes.add_textbox(Inches(0.5), top, col_w, height).text_frame
    left.word_wrap = True
    _write_column(left, s.get("left_title", ""), s.get("left_bullets", []))

    right = slide.shapes.add_textbox(Inches(0.5) + col_w, top, col_w, height).text_frame
    right.word_wrap = True
    _write_column(right, s.get("right_title", ""), s.get("right_bullets", []))

    _apply_notes(slide, s.get("notes"))


def _table_slide(prs: Presentation, s: dict[str, Any]) -> None:
    layout = prs.slide_layouts[SLIDE_LAYOUTS["table"]]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = s.get("title", "")
    headers = s.get("headers") or []
    rows = s.get("rows") or []
    if not headers and not rows:
        _apply_notes(slide, s.get("notes"))
        return
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = 1 + len(rows) if headers else len(rows)
    left = Inches(0.5)
    top = Inches(1.75)
    width = prs.slide_width - Inches(1.0)
    height = Inches(0.4) * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    if headers:
        for c, val in enumerate(headers):
            tbl.cell(0, c).text = str(val)
        start = 1
    else:
        start = 0
    for r, row in enumerate(rows, start=start):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = str(val)
    _apply_notes(slide, s.get("notes"))


def _image_slide(prs: Presentation, s: dict[str, Any]) -> None:
    layout = prs.slide_layouts[SLIDE_LAYOUTS["image"]]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = s.get("title", "")
    image_path = s.get("image_path")
    if image_path:
        top = Inches(1.75)
        slide.shapes.add_picture(
            str(image_path), Inches(0.5), top,
            width=prs.slide_width - Inches(1.0),
        )
    _apply_notes(slide, s.get("notes"))


_builders = {
    "title": _title_slide,
    "section": _section_slide,
    "bullets": _bullets_slide,
    "two_column": _two_column_slide,
    "table": _table_slide,
    "image": _image_slide,
}


def _find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _write_column(tf, title: str, bullets: list[str]) -> None:
    tf.clear()
    header = tf.paragraphs[0]
    header.text = title
    header.font.bold = True
    header.font.size = Pt(18)
    for b in bullets:
        para = tf.add_paragraph()
        para.text = f"• {b}"
        para.font.size = Pt(14)


def _apply_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Generate a .pptx from a spec.")
    ap.add_argument("spec", help="Path to a JSON or YAML spec.")
    ap.add_argument("--output", "-o", default="deck.pptx", help="Output .pptx path.")
    ap.add_argument("--template", default=None,
                    help="Optional .pptx template (master) to apply.")
    args = ap.parse_args(argv)

    spec_path = Path(args.spec)
    if not spec_path.exists():
        print(f"error: spec not found: {spec_path}", file=sys.stderr)
        return 2

    template = Path(args.template) if args.template else None
    prs = build_deck(load_spec(spec_path), template=template)
    out = Path(args.output)
    prs.save(str(out))
    print(f"wrote {out} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
