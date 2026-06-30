"""Tests for generate.py — build a deck, re-read it, assert structure."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from generate import build_deck


def _sample_spec() -> dict:
    return {
        "title": "Test Deck",
        "slides": [
            {"type": "title", "title": "Test Deck", "subtitle": "Subtitle"},
            {"type": "section", "title": "Section A"},
            {"type": "bullets", "title": "Findings", "bullets": ["One", "Two", "Three"]},
            {
                "type": "two_column",
                "title": "Supply vs demand",
                "left_title": "Supply",
                "left_bullets": ["A", "B"],
                "right_title": "Demand",
                "right_bullets": ["C", "D"],
            },
            {
                "type": "table",
                "title": "Scorecard",
                "headers": ["Metric", "Value"],
                "rows": [["Rev", "100"], ["Cost", "40"]],
            },
        ],
    }


def test_builds_all_slide_types(tmp_path: Path) -> None:
    prs = build_deck(_sample_spec())
    assert len(prs.slides) == 5


def test_title_and_subtitle_written(tmp_path: Path) -> None:
    prs = build_deck(_sample_spec())
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    re = Presentation(str(out))
    assert re.slides[0].shapes.title.text == "Test Deck"


def test_bullets_render(tmp_path: Path) -> None:
    prs = build_deck(_sample_spec())
    bullets_slide = prs.slides[2]
    body = None
    for ph in bullets_slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body = ph
    assert body is not None
    text = body.text_frame.text
    assert "One" in text and "Two" in text and "Three" in text


def test_table_rows_rendered(tmp_path: Path) -> None:
    prs = build_deck(_sample_spec())
    table_slide = prs.slides[4]
    tables = [s for s in table_slide.shapes if s.has_table]
    assert tables, "expected a table on the table slide"
    tbl = tables[0].table
    assert tbl.cell(0, 0).text == "Metric"
    assert tbl.cell(1, 1).text == "100"


def test_unknown_slide_type_raises() -> None:
    with pytest.raises(ValueError):
        build_deck({"slides": [{"type": "sparkline"}]})


def test_notes_applied() -> None:
    prs = build_deck({
        "slides": [
            {"type": "section", "title": "Agenda", "notes": "Speaker notes here"},
        ]
    })
    assert prs.slides[0].notes_slide.notes_text_frame.text == "Speaker notes here"
