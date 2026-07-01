"""Tests for extract.py — fixtures are generated in setup, not committed.

Run:  pytest test_extract.py
Dependencies in requirements.txt.
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from extract import extract, UnsupportedFormatError


@pytest.fixture
def tmp_docx(tmp_path: Path) -> Path:
    from docx import Document
    doc = Document()
    doc.add_paragraph("Executive summary: placeholder content.")
    doc.add_paragraph("Second paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "h1"
    table.cell(0, 1).text = "h2"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    p = tmp_path / "sample.docx"
    doc.save(str(p))
    return p


@pytest.fixture
def tmp_xlsx(tmp_path: Path) -> Path:
    import pandas as pd
    p = tmp_path / "sample.xlsx"
    with pd.ExcelWriter(p, engine="openpyxl") as writer:
        pd.DataFrame({"a": [1, 2], "b": [3, 4]}).to_excel(writer, sheet_name="Summary", index=False)
        pd.DataFrame({"x": [10], "y": [20]}).to_excel(writer, sheet_name="Q3", index=False)
    return p


@pytest.fixture
def tmp_pdf(tmp_path: Path) -> Path:
    # pypdf can't easily write text PDFs; use reportlab if available, else skip.
    reportlab = pytest.importorskip("reportlab.pdfgen.canvas")
    p = tmp_path / "sample.pdf"
    c = reportlab.Canvas(str(p))
    c.drawString(72, 720, "Page one text for extraction.")
    c.showPage()
    c.drawString(72, 720, "Page two text.")
    c.showPage()
    c.save()
    return p


@pytest.fixture
def tmp_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title slide"
    slide.placeholders[1].text = "Subtitle body text"
    p = tmp_path / "sample.pptx"
    prs.save(str(p))
    return p


@pytest.fixture
def tmp_csv(tmp_path: Path) -> Path:
    p = tmp_path / "sample.csv"
    p.write_text("a,b\n1,2\n3,4\n")
    return p


@pytest.fixture
def tmp_json(tmp_path: Path) -> Path:
    p = tmp_path / "sample.json"
    p.write_text(json.dumps({"k": "v", "n": 1}))
    return p


@pytest.fixture
def tmp_yaml(tmp_path: Path) -> Path:
    p = tmp_path / "sample.yaml"
    p.write_text("k: v\nn: 1\n")
    return p


@pytest.fixture
def tmp_txt(tmp_path: Path) -> Path:
    p = tmp_path / "sample.txt"
    p.write_text("plain text\nsecond line\n")
    return p


def test_docx(tmp_docx: Path) -> None:
    e = extract(tmp_docx)
    assert e.type == "docx"
    assert any("Executive summary" in para for para in e.content["paragraphs"])
    assert e.content["tables"][0][0] == ["h1", "h2"]
    assert e.metadata["paragraph_count"] >= 2
    assert e.metadata["table_count"] == 1


def test_xlsx(tmp_xlsx: Path) -> None:
    e = extract(tmp_xlsx)
    assert e.type == "excel"
    assert set(e.content.keys()) == {"Summary", "Q3"}
    assert e.content["Summary"]["rows"] == [{"a": 1, "b": 3}, {"a": 2, "b": 4}]


def test_pdf(tmp_pdf: Path) -> None:
    e = extract(tmp_pdf)
    assert e.type == "pdf"
    assert "Page one" in e.content["pages"][0]
    assert e.metadata["page_count"] == 2


def test_pdf_max_pages(tmp_pdf: Path) -> None:
    e = extract(tmp_pdf, max_pages=1)
    assert e.metadata["pages_extracted"] == 1
    assert any("truncated" in flag for flag in e.flags)


def test_pptx(tmp_pptx: Path) -> None:
    e = extract(tmp_pptx)
    assert e.type == "pptx"
    assert e.content["slides"][0]["title"] == "Title slide"


def test_csv(tmp_csv: Path) -> None:
    e = extract(tmp_csv)
    assert e.type == "csv"
    assert e.content["rows"] == [{"a": 1, "b": 2}, {"a": 3, "b": 4}]


def test_json(tmp_json: Path) -> None:
    e = extract(tmp_json)
    assert e.type == "json"
    assert e.content == {"k": "v", "n": 1}


def test_yaml(tmp_yaml: Path) -> None:
    e = extract(tmp_yaml)
    assert e.type == "yaml"
    assert e.content == {"k": "v", "n": 1}


def test_txt(tmp_txt: Path) -> None:
    e = extract(tmp_txt)
    assert e.type == "text"
    assert "plain text" in e.content


def test_unsupported_extension(tmp_path: Path) -> None:
    p = tmp_path / "foo.unknown"
    p.write_text("x")
    with pytest.raises(UnsupportedFormatError):
        extract(p)


def test_missing_file(tmp_path: Path) -> None:
    with pytest.raises(FileNotFoundError):
        extract(tmp_path / "nope.docx")
