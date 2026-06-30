"""Extract structured content from enterprise document formats.

Usage:
    python extract.py <path> [--format json|text] [--max-pages N]

Supports .docx, .xlsx, .pdf, .pptx, .csv, .json, .yaml, .txt, .md
and returns structured content with provenance. Images (.png/.jpg) are
recognized but intended for multimodal-model consumption — extraction
falls back to OCR if pytesseract is available.

This script is kept small and dependency-light by design. Import what
you need; the installer hints in requirements.txt tell you what to add.
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class Extraction:
    type: str
    source: str
    method: str
    content: Any
    metadata: dict[str, Any] = field(default_factory=dict)
    flags: list[str] = field(default_factory=list)


class UnsupportedFormatError(ValueError):
    """Raised when the file extension is not in the dispatch table."""


def extract(path: str | Path, *, max_pages: int | None = None) -> Extraction:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(p)
    ext = p.suffix.lower()

    dispatch = {
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xlsx,
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
        ".csv": _extract_csv,
        ".json": _extract_json,
        ".yaml": _extract_yaml,
        ".yml": _extract_yaml,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".png": _extract_image,
        ".jpg": _extract_image,
        ".jpeg": _extract_image,
    }
    handler = dispatch.get(ext)
    if handler is None:
        raise UnsupportedFormatError(f"No handler for {ext}")
    return handler(p, max_pages=max_pages)


def _extract_docx(p: Path, **_) -> Extraction:
    from docx import Document  # python-docx
    doc = Document(str(p))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    tables = [
        [[cell.text for cell in row.cells] for row in tbl.rows]
        for tbl in doc.tables
    ]
    return Extraction(
        type="docx",
        source=str(p),
        method="python-docx",
        content={"paragraphs": paragraphs, "tables": tables},
        metadata={"paragraph_count": len(paragraphs), "table_count": len(tables)},
    )


def _extract_xlsx(p: Path, **_) -> Extraction:
    import pandas as pd  # requires openpyxl for .xlsx
    sheets = pd.read_excel(p, sheet_name=None, engine="openpyxl")
    return Extraction(
        type="excel",
        source=str(p),
        method="pandas+openpyxl",
        content={
            name: {
                "columns": list(df.columns),
                "rows": df.to_dict(orient="records"),
                "shape": list(df.shape),
            }
            for name, df in sheets.items()
        },
        metadata={"sheet_count": len(sheets), "sheet_names": list(sheets.keys())},
    )


def _extract_pdf(p: Path, *, max_pages: int | None = None, **_) -> Extraction:
    from pypdf import PdfReader
    reader = PdfReader(str(p))
    total = len(reader.pages)
    limit = min(total, max_pages) if max_pages else total
    pages = [reader.pages[i].extract_text() or "" for i in range(limit)]
    flags: list[str] = []
    if all(not page.strip() for page in pages):
        flags.append("pdf_appears_scanned_or_empty")
    if max_pages and total > max_pages:
        flags.append(f"truncated_to_first_{max_pages}_of_{total}_pages")
    return Extraction(
        type="pdf",
        source=str(p),
        method="pypdf",
        content={"pages": pages},
        metadata={"page_count": total, "pages_extracted": limit},
        flags=flags,
    )


def _extract_pptx(p: Path, **_) -> Extraction:
    from pptx import Presentation
    prs = Presentation(str(p))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        body: list[str] = []
        notes = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if shape == slide.shapes.title and title is None:
                        title = text
                    else:
                        body.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"index": idx, "title": title, "body": body, "notes": notes})
    return Extraction(
        type="pptx",
        source=str(p),
        method="python-pptx",
        content={"slides": slides},
        metadata={"slide_count": len(slides)},
    )


def _extract_csv(p: Path, **_) -> Extraction:
    import pandas as pd
    df = pd.read_csv(p)
    return Extraction(
        type="csv",
        source=str(p),
        method="pandas",
        content={"columns": list(df.columns), "rows": df.to_dict(orient="records")},
        metadata={"shape": list(df.shape)},
    )


def _extract_json(p: Path, **_) -> Extraction:
    with p.open() as fh:
        data = json.load(fh)
    return Extraction(type="json", source=str(p), method="stdlib", content=data)


def _extract_yaml(p: Path, **_) -> Extraction:
    import yaml  # PyYAML
    with p.open() as fh:
        data = yaml.safe_load(fh)
    return Extraction(type="yaml", source=str(p), method="PyYAML", content=data)


def _extract_text(p: Path, **_) -> Extraction:
    text = p.read_text(encoding="utf-8", errors="replace")
    return Extraction(
        type="text",
        source=str(p),
        method="read_text",
        content=text,
        metadata={"char_count": len(text), "line_count": text.count("\n") + 1},
    )


def _extract_image(p: Path, **_) -> Extraction:
    try:
        import pytesseract
        from PIL import Image
        text = pytesseract.image_to_string(Image.open(p))
        return Extraction(
            type="image",
            source=str(p),
            method="pytesseract",
            content=text,
            metadata={"char_count": len(text)},
        )
    except ImportError:
        return Extraction(
            type="image",
            source=str(p),
            method="none",
            content=None,
            flags=["use_multimodal_read_directly_or_install_pytesseract"],
        )


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Extract content from enterprise documents.")
    ap.add_argument("path", help="Path to the document.")
    ap.add_argument("--format", choices=("json", "text"), default="json")
    ap.add_argument("--max-pages", type=int, default=None,
                    help="For PDFs: limit pages extracted.")
    args = ap.parse_args(argv)

    try:
        result = extract(args.path, max_pages=args.max_pages)
    except UnsupportedFormatError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    except FileNotFoundError as exc:
        print(f"error: file not found: {exc}", file=sys.stderr)
        return 2

    if args.format == "json":
        print(json.dumps(asdict(result), indent=2, default=str))
    else:
        print(f"# {result.type} @ {result.source} ({result.method})")
        print(result.content if isinstance(result.content, str) else json.dumps(result.content, indent=2, default=str))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
